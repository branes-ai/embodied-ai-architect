#!/usr/bin/env python3
"""Build the Dependency Trap pitch deck as .pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# -- Color palette --
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DARK = RGBColor(0x16, 0x21, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT = RGBColor(0xE8, 0x4D, 0x3D)  # red-orange for emphasis
ACCENT2 = RGBColor(0x3D, 0x9B, 0xE9)  # blue for secondary
MUTED = RGBColor(0x88, 0x88, 0x99)
TABLE_BG = RGBColor(0x22, 0x2B, 0x45)
TABLE_HEADER = RGBColor(0x2A, 0x3A, 0x5C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def add_dark_bg(slide):
    """Fill slide with dark background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with a single run of text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return tf


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=20, color=LIGHT_GRAY, spacing=Pt(12)):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        run = p.add_run()
        run.text = f"  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
        # bullet dash
        bullet_run = p.add_run()
        bullet_run.text = ""
        p.text = ""
        run2 = p.add_run()
        run2.text = f"\u2014  {item}"
        run2.font.size = Pt(font_size)
        run2.font.color.rgb = color
        run2.font.name = "Calibri"
    return tf


def add_accent_line(slide, left, top, width):
    """Add a thin accent-colored horizontal line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


# ============================================================
# SLIDE 1: The Dependency Trap
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_dark_bg(slide1)

add_text_box(slide1, Inches(1), Inches(1.2), Inches(11), Inches(1.2),
             "YOUR AUTONOMY RUNS ON\nSOMEONE ELSE'S ROADMAP",
             font_size=44, color=WHITE, bold=True)

add_accent_line(slide1, Inches(1), Inches(3.0), Inches(2))

bullets = [
    "Their highest-volume market isn't yours",
    "Allocation shortages strand programs overnight",
    "Export controls redraw the map without warning",
]
add_bullet_list(slide1, Inches(1), Inches(3.5), Inches(10), Inches(3),
                bullets, font_size=24, color=LIGHT_GRAY)

add_text_box(slide1, Inches(1), Inches(6.2), Inches(11), Inches(0.8),
             "If you don't own the silicon strategy, you don't own the autonomy.",
             font_size=20, color=ACCENT, bold=False)


# ============================================================
# SLIDE 2: The Gap Is Getting Worse
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide2)

add_text_box(slide2, Inches(1), Inches(0.8), Inches(11), Inches(1),
             "THE GAP IS GETTING WORSE",
             font_size=44, color=WHITE, bold=True)

add_accent_line(slide2, Inches(1), Inches(2.2), Inches(2))

# Three columns
col_data = [
    ("PHYSICS", ACCENT, "5W drone\n30W robot",
     "Workloads growing\nScaling slowing"),
    ("GEOPOLITICS", ACCENT2, "2 countries fab\nadvanced silicon",
     "Export controls\nexpanding yearly"),
    ("ECONOMICS", RGBColor(0x4E, 0xC9, 0xB0), "$50 chip creates\n$500 system penalty",
     "Competitors won't\noverpay forever"),
]

for i, (title, accent_color, constraint, trend) in enumerate(col_data):
    x = Inches(1 + i * 3.8)

    add_text_box(slide2, x, Inches(2.8), Inches(3.2), Inches(0.6),
                 title, font_size=22, color=accent_color, bold=True)

    add_text_box(slide2, x, Inches(3.5), Inches(3.2), Inches(1.2),
                 constraint, font_size=20, color=WHITE)

    add_text_box(slide2, x, Inches(5.0), Inches(3.2), Inches(1.2),
                 trend, font_size=18, color=MUTED)


# ============================================================
# SLIDE 3: The Solution
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide3)

add_text_box(slide3, Inches(1), Inches(0.8), Inches(11), Inches(1),
             "DESCRIBE THE MISSION.\nGET THE SILICON.",
             font_size=44, color=WHITE, bold=True)

add_accent_line(slide3, Inches(1), Inches(2.6), Inches(2))

# Left side: input
add_text_box(slide3, Inches(1), Inches(3.2), Inches(4.5), Inches(0.5),
             "INPUT", font_size=16, color=MUTED, bold=True)

add_text_box(slide3, Inches(1), Inches(3.7), Inches(4.5), Inches(1.5),
             '"Delivery drone, visual SLAM +\n detection, 30fps, under 5 watts"',
             font_size=22, color=LIGHT_GRAY)

# Arrow
add_text_box(slide3, Inches(5.8), Inches(4.0), Inches(1.5), Inches(1),
             "\u25B6", font_size=48, color=ACCENT, alignment=PP_ALIGN.CENTER)

# Right side: output
add_text_box(slide3, Inches(7.5), Inches(3.2), Inches(5), Inches(0.5),
             "OUTPUT", font_size=16, color=MUTED, bold=True)

add_text_box(slide3, Inches(7.5), Inches(3.7), Inches(5), Inches(1.5),
             "Validated SoC architecture\nwith synthesizable Verilog",
             font_size=22, color=WHITE, bold=True)

# Pipeline steps
steps = ["Workload Analysis", "Architecture Exploration",
         "Constraint Validation", "RTL Generation"]
for i, step in enumerate(steps):
    x = Inches(1 + i * 3)
    add_text_box(slide3, x, Inches(5.6), Inches(2.8), Inches(0.5),
                 step, font_size=16, color=ACCENT2)
    if i < len(steps) - 1:
        add_text_box(slide3, x + Inches(2.5), Inches(5.55), Inches(0.5), Inches(0.5),
                     "\u2192", font_size=20, color=MUTED)

add_text_box(slide3, Inches(1), Inches(6.4), Inches(11), Inches(0.6),
             "One engineer.  One afternoon.",
             font_size=22, color=ACCENT, bold=True)


# ============================================================
# SLIDE 4: Proof
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide4)

add_text_box(slide4, Inches(1), Inches(0.8), Inches(11), Inches(1),
             "IT WORKS TODAY",
             font_size=44, color=WHITE, bold=True)

add_accent_line(slide4, Inches(1), Inches(2.0), Inches(2))

# Three proof point cards
proofs = [
    ("DRONE SoC", "Auto-optimized to\nall-PASS constraints", "1.5 sec"),
    ("QUADRUPED SoC", "4 concurrent workloads\nPareto-ranked", "< 2 sec"),
    ("RTL PIPELINE", "Synthesizable Verilog\nverified with Yosys", "seconds"),
]

for i, (title, desc, time) in enumerate(proofs):
    x = Inches(1 + i * 3.8)
    # Card background
    card = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.6), Inches(3.3), Inches(3.2)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = TABLE_BG
    card.line.fill.background()

    add_text_box(slide4, x + Inches(0.3), Inches(2.9), Inches(2.7), Inches(0.5),
                 title, font_size=20, color=ACCENT, bold=True)

    add_text_box(slide4, x + Inches(0.3), Inches(3.6), Inches(2.7), Inches(1.2),
                 desc, font_size=20, color=WHITE)

    add_text_box(slide4, x + Inches(0.3), Inches(4.9), Inches(2.7), Inches(0.5),
                 time, font_size=28, color=ACCENT2, bold=True)

add_text_box(slide4, Inches(1), Inches(6.2), Inches(11), Inches(0.6),
             "50+ COTS platforms profiled  \u00B7  Custom accelerator RTL  \u00B7  No cloud dependency",
             font_size=18, color=MUTED)


# ============================================================
# SLIDE 5: The Ask
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide5)

add_text_box(slide5, Inches(1), Inches(0.8), Inches(11), Inches(1),
             "OWN THE COMPUTE.\nOWN THE AUTONOMY.",
             font_size=44, color=WHITE, bold=True)

add_accent_line(slide5, Inches(1), Inches(2.6), Inches(2))

# Three phase boxes
phases = [
    ("EVALUATE", "2 weeks", "Your workloads on\nCOTS + custom targets"),
    ("DESIGN", "8 weeks", "Validated SoC\narchitecture with RTL"),
    ("VERIFY", "12 weeks", "Pre-silicon validation\ntape-out ready"),
]

for i, (phase, timeline, desc) in enumerate(phases):
    x = Inches(1 + i * 3.8)
    # Phase box
    box = slide5.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.0), Inches(3.3), Inches(2.5)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = TABLE_BG
    box.line.fill.background()

    add_text_box(slide5, x + Inches(0.3), Inches(3.2), Inches(2.7), Inches(0.5),
                 phase, font_size=20, color=ACCENT2, bold=True)

    add_text_box(slide5, x + Inches(0.3), Inches(3.8), Inches(2.7), Inches(0.6),
                 timeline, font_size=32, color=WHITE, bold=True)

    add_text_box(slide5, x + Inches(0.3), Inches(4.5), Inches(2.7), Inches(1),
                 desc, font_size=18, color=LIGHT_GRAY)

    # Arrow between boxes
    if i < len(phases) - 1:
        add_text_box(slide5, x + Inches(3.3), Inches(3.9), Inches(0.5), Inches(0.8),
                     "\u25B6", font_size=28, color=MUTED, alignment=PP_ALIGN.CENTER)

add_text_box(slide5, Inches(1), Inches(6.0), Inches(11), Inches(0.8),
             "Next step:  2-week evaluation sprint.\nYou bring models.  We bring the platform.",
             font_size=22, color=WHITE, bold=True)

# -- Save --
out_path = "docs/pitch-deck-dependency-trap.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
